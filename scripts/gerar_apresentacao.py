#!/usr/bin/env python3
"""Gera a apresentacao PPTX do Digital Twin Processo Bayer (16:9, paleta Midnight)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------------------------------
# Paleta (Midnight Execution + acentos do projeto)
# ----------------------------------------------------------------------------
DARK    = "0B0F19"
NAVY    = "1E2761"
ICE     = "CADCFC"
WHITE   = "FFFFFF"
OFF     = "F7F9FC"
ACCENT  = "38BDF8"
GREEN   = "10B981"
AMBER   = "F59E0B"
TEXT_D  = "111827"
MUTED   = "6B7280"
PANEL   = "111827"
BERRY   = "6D2E46"

HFONT = "Georgia"
BFONT = "Calibri"

SW, SH = 13.333, 7.5


def C(hexstr: str) -> RGBColor:
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))


def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs_width(), prs_height())
    shp.fill.solid(); shp.fill.fore_color.rgb = C(color)
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def prs_width():
    return Inches(SW)


def prs_height():
    return Inches(SH)


def rect(slide, l, t, w, h, fill, line=None, lw=1.0, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = C(fill)
    if line:
        shp.line.color.rgb = C(line); shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def text(slide, l, t, w, h, txt, size=18, color=TEXT_D, bold=False, italic=False,
         align=PP_ALIGN.LEFT, font=BFONT, anchor=MSO_ANCHOR.TOP, spacing=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.autosize = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    lines = str(txt).split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = C(color); r.font.name = font
    return tb


def header(slide, kicker, title):
    text(slide, 0.7, 0.4, 11.8, 0.4, kicker, size=12, color=ACCENT,
         bold=True, font=HFONT)
    text(slide, 0.7, 0.82, 11.8, 0.85, title, size=30, color=NAVY,
         bold=True, font=HFONT)
    rect(slide, 0.72, 1.72, 1.6, 0.09, ACCENT)


def footer(slide, n):
    text(slide, 0.7, 7.08, 9.0, 0.3,
         "Digital Twin Processo Bayer  ·  projeto_bayer", size=9, color=MUTED)
    text(slide, 12.1, 7.08, 1.0, 0.3, str(n), size=10, color=MUTED,
         align=PP_ALIGN.RIGHT)


def bullets(slide, l, t, w, h, items, size=15, color=TEXT_D, gap_after=6):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.autosize = MSO_AUTO_SIZE.NONE
    first = True
    for head, body in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap_after)
        r = p.add_run(); r.text = "▸ "
        r.font.color.rgb = C(ACCENT); r.font.size = Pt(size); r.font.bold = True
        if head:
            rh = p.add_run(); rh.text = head
            rh.font.size = Pt(size); rh.font.bold = True; rh.font.color.rgb = C(color)
        if body:
            rb = p.add_run(); rb.text = (" — " if head else " ") + body
            rb.font.size = Pt(size); rb.font.color.rgb = C(TEXT_D)
    return tb


def chip(slide, n, label, l, t, w, h, accent=ACCENT):
    rect(slide, l, t, 0.42, 0.42, accent, rounded=True)
    text(slide, l, t + 0.02, 0.42, 0.38, str(n), size=16, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, l + 0.6, t - 0.04, w - 0.6, h, label, size=14, color=TEXT_D,
         anchor=MSO_ANCHOR.MIDDLE)


# ============================================================================
prs = new_prs()

# ---------------- 1. CAPA ----------------
s = blank(prs); bg(s, DARK)
rect(s, 0, 0, 0.38, SH, ACCENT)
rect(s, 0.38, 0, 0.05, SH, BERRY)
text(s, 1.1, 2.05, 11.0, 0.45, "CONTROLE PREDITIVO · IA APLICADA À ENG. QUÍMICA",
     size=13, color=ACCENT, bold=True, font=HFONT)
text(s, 1.1, 2.5, 11.5, 1.9, "Digital Twin do\nProcesso Bayer", size=50,
     color=WHITE, bold=True, font=HFONT)
text(s, 1.1, 4.45, 11.2, 0.9,
     "Sistema de supervisão e controle com LangGraph + Fuzzy Adaptativo\n+ Human-in-the-Loop (HITL)",
     size=17, color=ICE)
rect(s, 1.1, 6.05, 3.2, 0.05, GREEN)
text(s, 1.1, 6.35, 11.5, 0.35,
     "projeto_bayer   ·   github.com/LuisGSVasconcelos/bayer-digital-twin   ·   v1.0.0   ·   agosto/2026",
     size=11, color=MUTED)

# ---------------- 2. AGENDA ----------------
s = blank(prs); bg(s, OFF)
header(s, "ROTEIRO", "Agenda")
agenda = [
    "Contexto e objetivo",
    "Topologia do processo (Bayer)",
    "Arquitetura do sistema",
    "Controlador Fuzzy Adaptativo",
    "Gêmeo Digital: distúrbios simulados",
    "Agente LangGraph + HITL",
    "Dashboard de supervisão",
    "Validação e métricas",
    "Conclusão e próximos passos",
]
y = 2.15
for i, item in enumerate(agenda, start=1):
    chip(s, i, item, 0.9, y, 11.0, 0.42)
    y += 0.55
footer(s, 2)

# ---------------- 3. CONTEXTO & OBJETIVO ----------------
s = blank(prs); bg(s, OFF)
header(s, "O PROBLEMA", "Contexto e objetivo")
text(s, 0.7, 2.2, 11.9, 1.3,
     "O projeto modela um Gêmeo Digital do processo de produção de alumina "
     "(Processo Bayer), focado na etapa crítica de controle de nível dos "
     "decantadores, para evitar transbordamentos — especialmente sob chuva intensa.",
     size=16, color=TEXT_D)
rect(s, 0.7, 3.6, 11.9, 0.06, ICE)
bullets(s, 0.7, 3.95, 6.95, 2.4, [
    ("Supervisão inteligente", "de nível dos decantadores (PA/PB) em série/paralelo"),
    ("Ambiente de testes (PoC)", "para validar estratégias de controle antes da planta real"),
    ("Decisões determinísticas + IA", "controlador Fuzzy adaptativo orquestrado por LangGraph"),
], size=15, gap_after=14)
rect(s, 8.0, 3.95, 4.6, 2.55, NAVY, rounded=True)
text(s, 8.35, 4.25, 3.9, 0.4, "DIFERENCIAIS", size=12, color=ACCENT, bold=True, font=HFONT)
text(s, 8.35, 4.7, 3.9, 0.35, "HITL obrigatório", size=16, color=WHITE, bold=True)
text(s, 8.35, 5.02, 3.9, 0.6, "ação física só após aprovação do operador", size=11, color=ICE)
text(s, 8.35, 5.65, 3.9, 0.35, "Roda offline", size=16, color=WHITE, bold=True)
text(s, 8.35, 5.97, 3.9, 0.4, "núcleo testado sem serviços externos", size=11, color=ICE)
footer(s, 3)

# ---------------- 4. TOPOLOGIA ----------------
s = blank(prs); bg(s, OFF)
header(s, "MODELO FÍSICO", "Topologia do processo (simplificação Bayer)")
steps = [
    ("S1", "Digestão", NAVY),
    ("S2", "Digestão", NAVY),
    ("Divisão", "fluxo", "6B7280"),
    ("PA", "Decantador", AMBER),
    ("PB", "Decantador", AMBER),
]
x = 0.7
for name, cap, col in steps:
    rect(s, x, 3.15, 1.95, 1.25, col, rounded=True)
    text(s, x, 3.42, 1.95, 0.5, name, size=20, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
    text(s, x, 3.9, 1.95, 0.4, cap, size=11, color=WHITE, align=PP_ALIGN.CENTER)
    x += 1.95
    if name != "PB":
        text(s, x - 0.28, 3.42, 0.5, 0.5, "→", size=22, color=ACCENT,
             bold=True, align=PP_ALIGN.CENTER)
        x += 0.55
text(s, 0.7, 4.75, 11.9, 1.6,
     "Carga (bauxita + soda) atravessa os reatores de digestão em SÉRIE e é dividida entre "
     "os decantadores em PARALELO. Os decantadores são abertos e acumulam chuva na área "
     "exposta (150 m²) — daí o risco de transbordo sob temporal.",
     size=14, color=TEXT_D).text_frame.paragraphs[0].line_spacing = 1.2
rect(s, 0.7, 6.15, 11.9, 0.7, ICE, rounded=True)
text(s, 0.9, 6.23, 5.6, 0.55, "SÉRIE — reatores de digestão (fechados)", size=12, color=NAVY, bold=True)
text(s, 7.0, 6.23, 5.6, 0.55, "PARALELO — decantadores (abertos, 150 m²)", size=12, color=NAVY, bold=True)
footer(s, 4)

# ---------------- 5. ARQUITETURA ----------------
s = blank(prs); bg(s, WHITE)
header(s, "VISÃO DE SISTEMA", "Arquitetura do sistema")
bullets(s, 0.6, 2.3, 3.9, 4.3, [
    ("LangGraph", "grafo de estados com ciclo"),
    ("Fuzzy Adaptativo", "por decantador (PA/PB)"),
    ("HITL", "interrupt_before p/ aprovação"),
    ("Persistência", "InfluxDB (time-series)"),
    ("Dashboard", "Streamlit + Plotly"),
], size=13, gap_after=12)
img = s.shapes.add_picture("assets/arquitetura_processo_bayer.png",
                           Inches(4.75), Inches(1.95), width=Inches(8.05))
footer(s, 5)

# ---------------- 6. CONTROLADOR FUZZY ADAPTATIVO ----------------
s = blank(prs); bg(s, OFF)
header(s, "INTELIGÊNCIA DE CONTROLE", "Controlador Fuzzy Adaptativo")
rect(s, 0.7, 2.2, 6.0, 3.4, WHITE, line=ICE, rounded=True)
text(s, 1.0, 2.35, 5.3, 0.4, "Base Fuzzy — Mamdani + Centroide", size=14, color=NAVY, bold=True)
bullets(s, 1.0, 2.85, 5.4, 2.6, [
    ("Entradas", "erro (nível − setpoint) e derivada"),
    ("5×5 regras", "set trapezoidais (NG…PG)"),
    ("Saída", "abertura da válvula em [0,1]"),
    ("Domínio do erro", "±20% → até erro +35 (tanque cheio)"),
], size=12, gap_after=9)
rect(s, 6.95, 2.2, 5.7, 3.4, NAVY, rounded=True)
text(s, 7.25, 2.35, 5.0, 0.4, "Adaptativo — ganho multiplicador γ", size=14, color=ACCENT, bold=True)
text(s, 7.25, 2.9, 5.0, 0.7, "u_real = γ · u_fuzzy", size=20, color=WHITE, bold=True, font="Consolas")
text(s, 7.25, 3.6, 5.1, 0.4, "Δγ = η·e·u + α·Δγ_prev", size=15, color=ICE, font="Consolas")
text(s, 7.25, 4.1, 5.1, 1.3,
     "Gradiente descendente com momentum + saturação γ ∈ [0.3, 2.5]. "
     "Aprende a agressividade da ação, com controlador independente por decantador.",
     size=12, color=ICE)
rect(s, 0.7, 5.85, 11.9, 0.95, DARK, rounded=True)
text(s, 1.0, 6.0, 5.6, 0.6, "Enquanto o fuzzy base define as regras, o ganho γ\nadapta a escala ao longo do tempo.",
     size=12, color=WHITE, bold=True)
text(s, 7.2, 6.12, 5.2, 0.5, "Melhoria reproduzível de RMSE no benchmark.", size=12, color=ACCENT)
footer(s, 6)

# ---------------- 7. DISTÚRBIOS ----------------
s = blank(prs); bg(s, OFF)
header(s, "GÊMEO DIGITAL", "Distúrbios operacionais simulados")
disturbios = [
    ("Variação de alimentação", "oscilação + ruído na vazão de entrada"),
    ("Desgaste de bomba", "fouling — redução gradual da vazão"),
    ("Stiction de válvula", "atuador não responde instantaneamente"),
    ("Desbalanço de paralelo", "divisão PA/PB varia (20–80%)"),
    ("Spike de sensor", "pico até +15% no nível medido"),
    ("Diluição de TC", "água adicionada ao teor cáustico"),
    ("Variação de sílica", "teor de SiO₂ oscila na bauxita"),
]
cw, ch, gap = 3.83, 1.7, 0.2
x0, y0 = 0.7, 2.35
for i, (name, desc) in enumerate(disturbios):
    col = i % 3
    row = i // 3
    x = x0 + col * (cw + gap)
    y = y0 + row * (ch + gap)
    rect(s, x, y, cw, ch, WHITE, line=ICE, rounded=True)
    text(s, x + 0.25, y + 0.18, cw - 0.5, 0.4, name, size=13.5, color=NAVY, bold=True)
    text(s, x + 0.25, y + 0.62, cw - 0.5, 0.9, desc, size=11, color=TEXT_D)
footer(s, 7)

# ---------------- 8. AGENTE LANGGRAPH + HITL ----------------
s = blank(prs); bg(s, OFF)
header(s, "ORQUESTRAÇÃO", "Agente LangGraph + Human-in-the-Loop")
flow = [("coleta", ACCENT), ("analise", ACCENT),
        ("calcular", ACCENT), ("aguardar", AMBER), ("executar", ACCENT)]
fx, fw, fh = 0.7, 2.0, 1.2
for i, (name, col) in enumerate(flow):
    rect(s, fx, 2.7, fw, fh, col, rounded=True)
    text(s, fx, 2.98, fw, 0.5, name, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    fx += fw
    if i != len(flow) - 1:
        text(s, fx - 0.28, 2.92, 0.5, 0.5, "→", size=20, color=TEXT_D, bold=True, align=PP_ALIGN.CENTER)
        fx += 0.55
text(s, 0.7, 4.15, 11.9, 1.6,
     "coleta → análise → calcular_controle → [risco crítico] aguardar_operador → executar_controle → END\n\n"
     "O fluxo interrompe (interrupt_before) antes da ação física: abertura de válvula só ocorre após "
     "aprovação do operador, registrada via checkpoint (MemorySaver).",
     size=13, color=TEXT_D).text_frame.paragraphs[1].line_spacing = 1.15
rect(s, 0.7, 5.9, 11.9, 0.9, AMBER, rounded=True)
text(s, 1.0, 6.0, 11.3, 0.7,
     "⚠  Segurança: o LLM/IA não decide regras físicas. Criticidade é determinística (nível filtrado por EMA + chuva) "
     "e a validação humana é obrigatória.",
     size=12.5, color=TEXT_D, bold=True).text_frame.paragraphs[0].line_spacing = 1.1
footer(s, 8)

# ---------------- 9. DASHBOARD ----------------
s = blank(prs); bg(s, OFF)
header(s, "SUPERVISÃO EM TEMPO REAL", "Dashboard interativo (Streamlit + Plotly)")
bullets(s, 0.7, 2.3, 6.4, 4.2, [
    ("KPIs", "níveis PA/PB, TC, perda de soda"),
    ("Gráficos", "PV × SP × MV, química, distúrbios"),
    ("Controles", "iniciar/parar e velocidade da simulação"),
    ("HITL no painel", "botão de aprovação de ação emergencial"),
    ("Log de eventos", "últimos ciclos com alertas e variáveis"),
], size=14, gap_after=13)
# mock de janela com barras decorativas
rect(s, 7.3, 2.3, 5.35, 4.2, WHITE, line=ICE, rounded=True)
rect(s, 7.3, 2.3, 5.35, 0.5, NAVY, rounded=False)
text(s, 7.5, 2.36, 5.0, 0.4, "● ● ●   Sala de Controle", size=11, color=WHITE)
rect(s, 7.55, 3.0, 2.2, 0.5, ACCENT, rounded=True)
text(s, 7.55, 3.05, 2.2, 0.4, "▶ Iniciar", size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
rect(s, 7.55, 3.7, 4.85, 0.55, OFF, line="D6DEEA", rounded=True)
rect(s, 7.75, 3.85, 1.2, 0.25, GREEN)
rect(s, 9.2, 3.85, 0.9, 0.25, AMBER)
rect(s, 10.3, 3.85, 0.7, 0.25, "D6DEEA")
rect(s, 7.55, 4.5, 4.85, 0.55, OFF, line="D6DEEA", rounded=True)
rect(s, 7.75, 4.6, 0.8, 0.3, NAVY)
rect(s, 8.7, 4.6, 1.5, 0.3, ACCENT)
rect(s, 10.35, 4.6, 1.2, 0.3, AMBER)
text(s, 7.55, 5.3, 4.9, 0.5, "Gráfico: Nível PA × Setpoint × Abertura", size=11, color=MUTED, italic=True)
footer(s, 9)

# ---------------- 10. VALIDAÇÃO & MÉTRICAS ----------------
s = blank(prs); bg(s, DARK)
text(s, 0.7, 0.7, 11.8, 0.4, "RESULTADOS", size=12, color=ACCENT, bold=True, font=HFONT)
text(s, 0.7, 1.1, 11.8, 0.85, "Validação e métricas", size=30, color=WHITE, bold=True, font=HFONT)
rect(s, 0.72, 2.0, 1.6, 0.09, ACCENT)
stats = [
    ("23", "testes automatizados"),
    ("8.14→7.81%", "RMSE (padrão → adaptativo)"),
    ("+4.1%", "melhoria reproduzível"),
    ("10", "runs com seeds fixas"),
]
x = 0.7
for big, lab in stats:
    rect(s, x, 2.6, 2.85, 2.0, PANEL, rounded=True)
    text(s, x, 2.85, 2.85, 0.9, big, size=30, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    text(s, x + 0.2, 3.85, 2.45, 0.7, lab, size=11, color=ICE, align=PP_ALIGN.CENTER)
    x += 3.05
bullets(s, 0.7, 5.1, 11.9, 1.6, [
    ("Métrica honesta", "benchmark com seed fixa e repetições (média ± desvio), não alegação de execução única"),
    ("Cobertura", "fuzzy, adaptativo, simulador físico e fluxo do agente (HITL)"),
    ("Validação executada", "suíte reproduzida no ambiente local antes da conclusão"),
], size=13.5, color=WHITE, gap_after=8)
footer(s, 10)

# ---------------- 11. CONCLUSÃO & PRÓXIMOS PASSOS ----------------
s = blank(prs); bg(s, DARK)
rect(s, 0, 0, 0.38, SH, ACCENT)
text(s, 1.1, 1.0, 11.0, 0.4, "ENCERRAMENTO", size=12, color=ACCENT, bold=True, font=HFONT)
text(s, 1.1, 1.4, 11.5, 0.9, "Conclusão e próximos passos", size=32, color=WHITE, bold=True, font=HFONT)
col_items = [
    ("ENTREGUE", [
        "Gêmeo Digital multitanque (série + paralelo)",
        "Controle Fuzzy Adaptativo por tanque",
        "Agente LangGraph com HITL obrigatório",
        "Dashboard Streamlit + 23 testes",
        "Manual (PDF) e repositório público",
    ]),
    ("PRÓXIMOS PASSOS", [
        "Ingestão de documentos (Docling) p/ calibrar o Gêmeo",
        "RAG no agente (explicação ao operador)",
        "Controle Preditivo (MPC)",
        "Integração com hardware (CLP/Modbus)",
    ]),
]
x = 1.1
for title, items in col_items:
    rect(s, x, 2.6, 5.5, 4.1, PANEL, rounded=True)
    text(s, x + 0.35, 2.85, 4.8, 0.4, title, size=14, color=ACCENT, bold=True, font=HFONT)
    ty = 3.35
    for it in items:
        text(s, x + 0.35, ty, 4.9, 0.6, "▸ " + it, size=12, color=ICE)
        ty += 0.62
    x += 6.1
text(s, 1.1, 6.9, 11.5, 0.4,
     "Disponível em github.com/LuisGSVasconcelos/bayer-digital-twin  ·  release v1.0.0 (manual PDF)",
     size=11, color=MUTED)

prs.save("Apresentacao_Projeto_Bayer.pptx")
print("OK: Apresentacao_Projeto_Bayer.pptx")