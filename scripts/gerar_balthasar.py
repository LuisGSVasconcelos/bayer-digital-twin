#!/usr/bin/env python3
"""Reconstrói a apresentação do Digital Twin usando o TEMPLATE Balthasar (SlidesCarnival).

Estratégia: abre o template como base (herda tema + layouts), remove os slides de
exemplo e recria cada lâmina a partir dos layouts do template, preenchendo os
placeholders (título/corpo) e adicionando shapes com as cores do tema.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy

TEMPLATE = "Balthasar · SlidesCarnival.pptx"
OUT = "Apresentacao_Projeto_Bayer_Balthasar.pptx"

# Cores do tema Balthasar
BLUE   = "058DC7"
GREEN  = "158158"
ORANGE = "ED561B"
INK    = "1F1F1F"
DARK   = "000000"
WHITE  = "FFFFFF"
GRAY   = "595959"
LITE   = "EAF3F2"

FONT = "Arial"
SW, SH = 10.0, 5.62


def C(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def get_layout(prs, name):
    for L in prs.slide_masters[0].slide_layouts:
        if L.name == name:
            return L
    return prs.slide_masters[0].slide_layouts[3]


def clear_slides(prs):
    """Remove todos os slides do template."""
    lst = prs.slides._sldIdLst
    for sld in list(lst):
        rId = sld.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.part.drop_rel(rId)
        lst.remove(sld)


def new_slide(prs, layout_name):
    return prs.slides.add_slide(get_layout(prs, layout_name))


def set_placeholder(slide, idx, lines, size=None, bold=False, color=None):
    """Preenche um placeholder com várias linhas (herdando estilo do layout)."""
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.clear()
    text = lines if isinstance(lines, list) else [lines]
    for i, ln in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln
        if size: r.font.size = Pt(size)
        if bold: r.font.bold = True
        if color: r.font.color.rgb = C(color)
    return ph


def add_text(slide, l, t, w, h, txt, size=14, color=INK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = str(txt).split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = C(color); r.font.name = FONT
    return tb


def add_rect(slide, l, t, w, h, fill, rounded=True, line=None, lw=1.0):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = C(fill)
    if line:
        shp.line.color.rgb = C(line); shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def footer(slide, n):
    add_text(slide, 0.45, 5.28, 7.0, 0.3,
             "Digital Twin Processo Bayer  ·  projeto_bayer", size=8, color=GRAY)
    add_text(slide, 9.0, 5.28, 0.6, 0.3, str(n), size=9, color=GRAY,
             align=PP_ALIGN.RIGHT)


prs = Presentation(TEMPLATE)
clear_slides(prs)

# ============================ S1 - CAPA ============================
s = new_slide(prs, "TITLE_1")
set_placeholder(s, 0, ["Digital Twin do", "Processo Bayer"])
set_placeholder(s, 1, ["Sistema de supervisão e controle · LangGraph + Fuzzy Adaptativo + HITL"])
add_text(s, 1.5, 4.7, 7.0, 0.3,
         "projeto_bayer · github.com/LuisGSVasconcelos/bayer-digital-twin · v1.0.0",
         size=10, color=GRAY, align=PP_ALIGN.CENTER)

# ============================ S2 - AGENDA ============================
s = new_slide(prs, "TITLE_AND_BODY")
set_placeholder(s, 0, "Agenda")
set_placeholder(s, 1, [
    "Contexto e objetivo",
    "Topologia do processo (Bayer)",
    "Arquitetura do sistema",
    "Controlador Fuzzy Adaptativo",
    "Gêmeo Digital: distúrbios simulados",
    "Agente LangGraph + HITL",
    "Dashboard de supervisão",
    "Validação e métricas",
    "Conclusão e próximos passos",
], size=14)
footer(s, 2)

# ============================ S3 - CONTEXTO ============================
s = new_slide(prs, "TITLE_AND_TWO_COLUMNS")
set_placeholder(s, 0, "Contexto e objetivo")
set_placeholder(s, 1, [
    "Gêmeo Digital da produção de alumina (Processo Bayer), focado no controle de "
    "nível dos decantadores para evitar transbordamentos sob chuva intensa.",
    "Supervisão inteligente — nível dos decantadores (PA/PB)",
    "Ambiente de testes (PoC) — valida antes da planta real",
    "Decisões determinísticas + IA — Fuzzy adaptativo + LangGraph",
], size=13)
set_placeholder(s, 2, [
    "HITL obrigatório — ação física só após aprovação do operador",
    "Roda offline — núcleo testado sem serviços externos",
    "Regras físicas de segurança não delegadas ao LLM",
], size=13)
footer(s, 3)

# ============================ S4 - TOPOLOGIA ============================
s = new_slide(prs, "TITLE_ONLY")
set_placeholder(s, 0, "Topologia do processo (simplificação Bayer)")
steps = [("S1", GREEN), ("S2", GREEN), ("Divisão", BLUE), ("PA", ORANGE), ("PB", ORANGE)]
W, GAP = 1.45, 0.44
x = 0.5
for i, (name, col) in enumerate(steps):
    add_rect(s, x, 2.0, W, 1.05, col)
    add_text(s, x, 2.2, W, 0.6, name, size=22, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
    if i != len(steps) - 1:
        add_text(s, x + W + 0.06, 2.15, 0.32, 0.6, "→", size=24, color=BLUE,
                 bold=True, align=PP_ALIGN.CENTER)
    x += W + GAP
add_text(s, 0.55, 3.35, 8.9, 1.0,
         "Carga percorre os reatores de digestão em SÉRIE e é dividida entre os "
         "decantadores em PARALELO. Os decantadores são abertos e captam chuva na "
         "área exposta (150 m²).", size=13, color=INK)
add_rect(s, 0.55, 4.5, 4.3, 0.62, LITE)
add_text(s, 0.75, 4.55, 4.0, 0.5, "SÉRIE — reatores fechados", size=12, color=GREEN, bold=True)
add_rect(s, 5.1, 4.5, 4.3, 0.62, LITE)
add_text(s, 5.3, 4.55, 4.0, 0.5, "PARALELO — decantadores abertos (150 m²)", size=12, color=ORANGE, bold=True)
footer(s, 4)

# ============================ S5 - ARQUITETURA ============================
s = new_slide(prs, "TITLE_ONLY")
set_placeholder(s, 0, "Arquitetura do sistema")
bullets = [
    ("LangGraph", " grafo de estados com ciclo"),
    ("Fuzzy Adaptativo", " por decantador (PA/PB)"),
    ("HITL", " interrupt_before p/ aprovação"),
    ("Persistência", " InfluxDB (time-series)"),
    ("Dashboard", " Streamlit + Plotly"),
]
add_text(s, 0.5, 2.0, 3.4, 2.6,
         "\n".join([h + " —" + b for h, b in bullets]), size=12, color=INK)
s.shapes.add_picture("assets/arquitetura_processo_bayer.png",
                     Inches(4.15), Inches(1.7), width=Inches(5.4))
footer(s, 5)

# ============================ S6 - FUZZY ============================
s = new_slide(prs, "TITLE_AND_TWO_COLUMNS")
set_placeholder(s, 0, "Controlador Fuzzy Adaptativo")
set_placeholder(s, 1, [
    "Base Fuzzy — Mamdani + Centroide",
    "Entradas: erro (nível − setpoint) e derivada",
    "Regras 5×5 com conjuntos trapezoidais",
    "Saída: abertura da válvula em [0,1]",
    "Domínio do erro até +35 (tanque cheio)",
], size=13)
set_placeholder(s, 2, [
    "Adaptativo — ganho multiplicador γ",
    "u_real = γ · u_fuzzy",
    "Δγ = η·e·u + α·Δγ_prev",
    "Saturação γ ∈ [0.3, 2.5]",
    "Controlador independente por tanque (PA/PB)",
], size=13)
footer(s, 6)

# ============================ S7 - DISTÚRBIOS ============================
s = new_slide(prs, "TITLE_AND_BODY")
set_placeholder(s, 0, "Gêmeo Digital — distúrbios operacionais simulados")
set_placeholder(s, 1, [
    "Variação de alimentação — oscilação + ruído na vazão de entrada",
    "Desgaste de bomba — fouling reduz gradualmente a vazão",
    "Stiction de válvula — atuador não responde instantaneamente",
    "Desbalanço de paralelo — divisão PA/PB varia (20–80%)",
    "Spike de sensor — pico até +15% no nível medido",
    "Diluição de TC — água adicionada ao teor cáustico",
    "Variação de sílica — teor de SiO₂ oscila na bauxita",
], size=14)
footer(s, 7)

# ============================ S8 - LANGGRAPH + HITL ============================
s = new_slide(prs, "TITLE_ONLY")
set_placeholder(s, 0, "Agente LangGraph + Human-in-the-Loop")
flow = [("coleta", BLUE), ("analise", BLUE), ("calcular", BLUE), ("aguardar", ORANGE), ("executar", BLUE)]
x = 0.5
for name, col in flow:
    add_rect(s, x, 2.05, 1.5, 0.9, col)
    add_text(s, x, 2.27, 1.5, 0.5, name, size=13, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
    x += 1.5
    if name != "executar":
        add_text(s, x - 0.26, 2.2, 0.52, 0.5, "→", size=22, color=DARK,
                 bold=True, align=PP_ALIGN.CENTER)
        x += 0.52
add_text(s, 0.5, 3.25, 9.0, 1.0,
         "coleta → análise → calcular_controle → [risco crítico] aguardar_operador → executar_controle → END\n"
         "O fluxo interrompe (interrupt_before) antes da ação física; a abertura de válvula "
         "exige aprovação do operador, registrada via checkpoint.", size=12, color=INK)
add_rect(s, 0.5, 4.45, 9.0, 0.7, "FCE4D6")
add_text(s, 0.75, 4.5, 8.5, 0.6,
         "⚠ Segurança: criticidade é determinística (nível filtrado por EMA + chuva) — "
         "o LLM/IA não decide regras físicas.", size=12, color=ORANGE, bold=True)
footer(s, 8)

# ============================ S9 - DASHBOARD ============================
s = new_slide(prs, "TITLE_AND_BODY")
set_placeholder(s, 0, "Dashboard interativo (Streamlit + Plotly)")
set_placeholder(s, 1, [
    "KPIs — níveis PA/PB, TC, perda de soda",
    "Gráficos — PV × SP × MV, química, distúrbios",
    "Controles — iniciar/parar e velocidade da simulação",
    "HITL no painel — botão de aprovação de ação emergencial",
    "Log de eventos — últimos ciclos com alertas e variáveis",
], size=14)
footer(s, 9)

# ============================ S10 - VALIDAÇÃO ============================
s = new_slide(prs, "TITLE_ONLY")
set_placeholder(s, 0, "Validação e métricas")
stats = [
    ("23", "testes automatizados", BLUE),
    ("8.14→7.81%", "RMSE padrão → adaptativo", GREEN),
    ("+4.1%", "melhoria reproduzível", ORANGE),
    ("10", "runs com seeds fixas", BLUE),
]
x = 0.55
for big, lab, col in stats:
    add_rect(s, x, 2.1, 2.0, 1.9, LITE)
    add_rect(s, x, 2.1, 0.12, 1.9, col, rounded=False)
    add_text(s, x + 0.2, 2.35, 1.7, 0.75, big, size=24, color=col, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, x + 0.15, 3.2, 1.75, 0.7, lab, size=10.5, color=INK,
             align=PP_ALIGN.CENTER)
    x += 2.25
add_text(s, 0.55, 4.35, 9.0, 0.8,
         "Benchmark com seed fixa e repetições (média ± desvio) — métrica honesta e "
         "reproduzível, cobrindo fuzzy, adaptativo, simulador físico e fluxo do agente (HITL).",
         size=12, color=INK)
footer(s, 10)

# ============================ S11 - CONCLUSÃO ============================
s = new_slide(prs, "TITLE_AND_TWO_COLUMNS")
set_placeholder(s, 0, "Conclusão e próximos passos")
set_placeholder(s, 1, [
    "Gêmeo Digital multitanque (série + paralelo)",
    "Controle Fuzzy Adaptativo por tanque",
    "Agente LangGraph com HITL obrigatório",
    "Dashboard Streamlit + 23 testes",
    "Manual (PDF) e repositório público",
], size=13)
set_placeholder(s, 2, [
    "Ingestão de documentos (Docling) p/ calibrar o Gêmeo",
    "RAG no agente (explicação ao operador)",
    "Controle Preditivo (MPC)",
    "Integração com hardware (CLP/Modbus)",
], size=13)
add_text(s, 0.55, 5.0, 9.0, 0.3,
         "github.com/LuisGSVasconcelos/bayer-digital-twin  ·  release v1.0.0 (manual PDF)",
         size=10, color=GRAY, align=PP_ALIGN.CENTER)

prs.save(OUT)
print("OK:", OUT)