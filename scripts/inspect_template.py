"""Inspeciona o template Balthasar: layouts, placeholders e slides."""
from pptx import Presentation
from pptx.util import Emu

PRS = Presentation("Balthasar · SlidesCarnival.pptx")
SW, SH = PRS.slide_width, PRS.slide_height
print(f"tamanho slide: {SW/914400:.2f} x {SH/914400:.2f} in")
print(f"n slides do master: {len(PRS.slide_masters[0].slide_layouts)}")

print("\n=== LAYOUTS do master ===")
for i, layout in enumerate(PRS.slide_masters[0].slide_layouts):
    phs = [(ph.placeholder_format.idx, ph.placeholder_format.type,
            (ph.name or '')) for ph in layout.placeholders]
    print(f"  L{i}: {layout.name!r}  placeholders={phs}")

for si, slide in enumerate(PRS.slides, start=1):
    ly = slide.slide_layout
    # layout index within the list
    try:
        li = PRS.slide_masters[0].slide_layouts.index(ly)
    except ValueError:
        li = "?"
    print(f"\n--- Slide {si} | layout={ly.name!r} (L{li}) ---")
    for ph in slide.placeholders:
        txt = (ph.text_frame.text.replace("\n", " / ")[:70] if ph.has_text_frame else "")
        print(f"   [ph{ph.placeholder_format.idx}:{ph.placeholder_format.type}] {ph.name!r}: {txt!r}")